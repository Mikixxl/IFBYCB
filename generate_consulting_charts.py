#!/usr/bin/env python3
"""
McKinsey / BCG Consulting Chart Generator
==========================================
Generates two files:
  • McKinsey_Charts.xlsx   – 12 chart sheets, data editable in-place
  • McKinsey_Slides.pptx   – 12 presentation slides (same charts)

Requirements:
  pip install xlsxwriter python-pptx

Usage:
  python generate_consulting_charts.py
"""

import sys

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BLUE   = "#002060"
MID_BLUE    = "#1F5FA6"
LIGHT_BLUE  = "#9DC3E6"
PALE_BLUE   = "#BDD7EE"
DARK_GRAY   = "#404040"
MID_GRAY    = "#7F7F7F"
LIGHT_GRAY  = "#D9D9D9"
GREEN       = "#70AD47"
RED         = "#C00000"
WHITE       = "#FFFFFF"
BLACK       = "#000000"

# ── Shared sample data ─────────────────────────────────────────────────────────

# Chart 1 – GBS Waterfall (stacked columns + floating reduction bars)
# Exact data from original McKinsey chart
GBS_CATS = [
    "Benchmark\nscope",
    "Process\nstandardization",
    "Strategic\nimportance",
    "System\nlimitations",
    "Signed-off\nscope",
    "Transferred\nto date",
    "To be\ntransferred",
]
# Milestone stacked segments (0 for reduction columns)
GBS_ACCOUNTING   = [89,  0,  0,  0, 66, 45, 21]
GBS_CUST_SERVICE = [45,  0,  0,  0, 31, 23,  8]
GBS_PROCUREMENT  = [34,  0,  0,  0, 19, 19,  5]
GBS_SALES_SUPP   = [25,  0,  0,  0, 13,  8,  0]
# Floating gray reduction bars: spacer positions the bar, reduction is the height
GBS_SPACER    = [  0, 172, 151, 129,   0,  0,  0]
GBS_REDUCTION = [  0,  21,  21,  22,   0,  0,  0]
GBS_TOTALS    = [193,   0,   0,   0, 129, 95, 34]  # for annotations

# Chart 2 – Diverging bar (Sentiment Survey)
SURVEY_STATEMENTS = [
    "Supply chain is well-optimised",
    "Supplier relationships are strong",
    "Procurement process is efficient",
    "Cost-savings targets are achievable",
    "Technology tools are adequate",
    "Sustainability goals are clear",
]
# [strongly_disagree%, disagree%, agree%, strongly_agree%, delta_pts]
SURVEY_DATA = [
    [-15, -22, 38, 25,  +8],
    [-10, -18, 42, 30,  -3],
    [-22, -28, 25, 25, +12],
    [ -5, -12, 50, 33,  +5],
    [-26, -19, 28, 27,  -7],
    [-12, -15, 40, 33,  +4],
]

# Chart 3 – Scatter / Line (GM Analysis)
PROMO_30    = [(100,45),(200,42),(320,39),(480,35),(650,31),(820,28),(950,25)]
BUY_2_GET_1 = [(150,40),(280,37),(410,33),(560,30),(710,27),(860,24),(980,22)]
GM_TARGET   = 30   # horizontal reference line

# Chart 4 – Doughnut (Revenue mix)
REV_CATS   = ["Electronics","Fashion","Home & Garden","Sports","Beauty","Other"]
REV_COLORS = [DARK_BLUE, MID_BLUE, LIGHT_BLUE, PALE_BLUE, MID_GRAY, LIGHT_GRAY]
REV_YEARS  = {
    "2019": [0.35, 0.20, 0.18, 0.12, 0.09, 0.06],
    "2021": [0.38, 0.22, 0.16, 0.11, 0.08, 0.05],
    "2024": [0.42, 0.19, 0.15, 0.10, 0.09, 0.05],
}

# Chart 6 – Mekko / Marimekko (Market landscape)
# Column width = market size (revenue $M); segment height = share within segment (%)
MEKKO_SEGMENTS = ["Europe", "North America", "Asia Pacific", "Rest of World"]
MEKKO_SEG_SIZES = [320, 480, 260, 140]   # market size $M (= column width)
MEKKO_SEG_COLORS = [DARK_BLUE, MID_BLUE, LIGHT_BLUE, PALE_BLUE]
# Row labels = product lines; values = % share within each segment (must sum to 100 per col)
MEKKO_ROWS   = ["Premium", "Mid-Range", "Value", "Other"]
MEKKO_ROW_COLORS = [DARK_BLUE, MID_BLUE, LIGHT_BLUE, LIGHT_GRAY]
MEKKO_DATA   = [           # share (%) per row per segment
    [40, 35, 25, 20],      # Premium
    [30, 30, 35, 30],      # Mid-Range
    [20, 25, 30, 35],      # Value
    [10, 10, 10, 15],      # Other
]
# Cumulative x-positions (for Mekko simulation via scatter/area approach)
# Total market = sum(MEKKO_SEG_SIZES) = 1200
MEKKO_TOTAL = sum(MEKKO_SEG_SIZES)

# Chart 5 – Range column + Line (Shipment volume)
SHIP_WEEKS  = ["Wk 1","Wk 2","Wk 3","Wk 4","Wk 5","Wk 6","Wk 7","Wk 8"]
SHIP_MIN    = [ 80,  85,  78,  90,  88,  82,  95,  91]
SHIP_MAX    = [120, 125, 115, 130, 128, 120, 140, 135]
SHIP_AVG    = [100, 105,  97, 110, 108, 101, 118, 113]
SHIP_ACTUAL = [ 95, 108, 102, 115, 112,  98, 125, 120]
SHIP_TARGET = [110]*8

# Chart 7 – 2×2 Priority Matrix
MATRIX_ITEMS = [
    # (label, ease_0to10, impact_0to10)
    ("Digital\nCore Platform",     8.0, 9.0),
    ("Customer 360",               6.5, 8.5),
    ("Cost Transformation",        7.5, 7.0),
    ("Talent Upskilling",          4.5, 6.0),
    ("Process Automation",         8.5, 6.5),
    ("Org Redesign",               3.0, 8.0),
    ("Analytics & AI",             5.5, 9.5),
    ("Supplier Consolidation",     7.0, 5.0),
    ("ESG Reporting",              6.0, 4.5),
    ("Legacy Migration",           2.5, 7.0),
]

# Chart 8 – EBITDA Bridge
# (label, value, kind)  kind: "start"|"up"|"down"|"total"
EBITDA_ITEMS = [
    ("Revenue",              500, "start"),
    ("Cost of Goods Sold",  -180, "down"),
    ("Gross Profit",         320, "total"),
    ("SG&A",                 -70, "down"),
    ("R&D",                  -45, "down"),
    ("D&A Add-back",         +32, "up"),
    ("One-off Items",        -20, "down"),
    ("EBITDA",               217, "total"),
]

# Chart 9 – Gantt / Transformation Roadmap
# (workstream, task, start_week, duration_weeks, color)
GANTT_DATA = [
    ("Assessment",  "Current State Analysis",   0,  6, "#002060"),
    ("Assessment",  "Stakeholder Interviews",    2,  5, "#002060"),
    ("Design",      "Future State Blueprint",    5,  8, "#1F5FA6"),
    ("Design",      "Business Case",             7,  6, "#1F5FA6"),
    ("Pilot",       "Site Preparation",         12,  4, "#9DC3E6"),
    ("Pilot",       "Pilot Execution",          14,  6, "#9DC3E6"),
    ("Rollout",     "Wave 1 Deployment",        19,  8, "#70AD47"),
    ("Rollout",     "Wave 2 Deployment",        25,  8, "#70AD47"),
    ("Rollout",     "Change Management",        19, 14, "#F39C12"),
    ("Governance",  "PMO & Tracking",            0, 32, "#D9D9D9"),
]
GANTT_TOTAL_WEEKS = 34

# Chart 10 – Harvey Balls (0=empty 1=quarter 2=half 3=three-quarter 4=full)
HARVEY_COMPANIES = ["McKinsey", "BCG", "Bain", "Deloitte", "Accenture"]
HARVEY_CRITERIA  = [
    "Strategy & Vision",
    "Digital Capability",
    "Operational Excellence",
    "M&A Advisory",
    "Talent & Culture",
    "Technology Delivery",
    "Sector Expertise",
]
HARVEY_SCORES = [
    [4, 4, 3, 3, 3],
    [3, 4, 2, 4, 4],
    [3, 3, 4, 3, 3],
    [4, 3, 3, 2, 2],
    [4, 3, 4, 3, 3],
    [2, 3, 2, 4, 4],
    [4, 3, 4, 2, 3],
]

# Chart 11 – Tornado / Sensitivity  (impact on project NPV, $M)
TORNADO_DRIVERS = [
    ("Revenue growth rate",        -48, +52),
    ("WACC / Discount rate",       -38, +30),
    ("EBITDA margin",              -35, +38),
    ("Exit multiple (EV/EBITDA)",  -32, +40),
    ("Working capital intensity",  -22, +18),
    ("Capex intensity",            -18, +15),
    ("Tax rate",                   -14, +12),
    ("Cost of goods sold",         -20, +16),
]

# Chart 12 – Sankey (Income Statement, values in $B)
SANKEY_SEGS = [
    ("Sales",                    2.1, "#26C6DA"),
    ("Service",                  2.3, "#EC407A"),
    ("Platform & Other",         1.8, "#AB47BC"),
    ("Marketing & Commerce",     1.3, "#FFA726"),
    ("Integration & Analytics",  1.3, "#42A5F5"),
]
SANKEY_TOTAL_REV    = 9.3
SANKEY_SUB_SUPPORT  = 8.8
SANKEY_PROF_SVCS    = 0.5
SANKEY_GROSS_PROFIT = 7.2
SANKEY_COR          = 2.1
SANKEY_OP_PROFIT    = 1.8
SANKEY_OP_EXP       = 5.4
SANKEY_NET_PROFIT   = 1.4
SANKEY_TAX          = 0.4
SANKEY_OPEX_DETAIL  = [
    ("Sales & Marketing",  3.2, "#EF5350"),
    ("R&D",                1.3, "#EF9A9A"),
    ("G&A",                0.7, "#FFCDD2"),
    ("Restructuring",      0.2, "#FF8A80"),
]


def _make_sankey_png():
    """Render Income Statement Sankey -> PNG BytesIO (requires matplotlib)."""
    import io
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches
        from matplotlib.path import Path
    except ImportError:
        return None

    W, H    = 16.0, 9.0
    fig     = plt.figure(figsize=(W, H), facecolor="white")
    ax      = fig.add_axes([0.0, 0.0, 1.0, 1.0])
    ax.set_xlim(0, W); ax.set_ylim(0, H); ax.axis("off")

    SCALE = 7.0 / SANKEY_TOTAL_REV
    Y0    = 0.75

    def vh(v): return v * SCALE

    def ribbon(x1, y1b, y1t, x2, y2b, y2t, color, alpha=0.45):
        mx = (x1 + x2) / 2
        verts = [(x1,y1b),(mx,y1b),(mx,y2b),(x2,y2b),(x2,y2t),(mx,y2t),(mx,y1t),(x1,y1t),(x1,y1b)]
        codes = [Path.MOVETO, Path.CURVE4,Path.CURVE4,Path.CURVE4,
                 Path.LINETO, Path.CURVE4,Path.CURVE4,Path.CURVE4, Path.CLOSEPOLY]
        ax.add_patch(mpatches.PathPatch(Path(verts,codes), facecolor=color, edgecolor="none", alpha=alpha))

    def nbar(x, yb, ht, color, w=0.20):
        ax.add_patch(plt.Rectangle((x-w/2,yb), w, ht, facecolor=color, edgecolor="white", linewidth=1.5, zorder=5))

    def lbl(x, y, text, color="#333", fs=9, ha="center", bold=False):
        ax.text(x, y, text, ha=ha, va="center", fontsize=fs, color=color,
                fontweight="bold" if bold else "normal")

    X = [1.0, 3.1, 5.2, 7.3, 10.1, 13.2]

    # Stage 0 - revenue segments
    y = Y0; seg_ys = []
    for _, v, _ in SANKEY_SEGS:
        seg_ys.append((y, y+vh(v))); y += vh(v)
    for (yb,yt),(_, _, c) in zip(seg_ys, SANKEY_SEGS):
        nbar(X[0], yb, yt-yb, c)
    for (yb,yt),(name,v,c) in zip(seg_ys, SANKEY_SEGS):
        lbl(X[0]-0.75,(yb+yt)/2,f"{name}\n${v}B",color=c,fs=8,ha="right")

    # Stage 1 - sub/support + prof services
    sub_b=Y0; sub_t=sub_b+vh(SANKEY_SUB_SUPPORT)
    pro_b=sub_t; pro_t=pro_b+vh(SANKEY_PROF_SVCS)
    y_dst=sub_b
    for (yb,yt),(_,v,c) in zip(seg_ys,SANKEY_SEGS):
        ribbon(X[0]+0.10,yb,yt,X[1]-0.10,y_dst,y_dst+vh(v),c,0.35); y_dst+=vh(v)
    nbar(X[1],sub_b,sub_t-sub_b,"#1565C0"); nbar(X[1],pro_b,pro_t-pro_b,"#283593")
    lbl(X[1]+0.65,(sub_b+sub_t)/2,f"Subscription\n& Support\n${SANKEY_SUB_SUPPORT}B",
        color="#0D47A1",fs=8,ha="left",bold=True)
    lbl(X[1]+0.65,(pro_b+pro_t)/2,f"Prof. Services\n${SANKEY_PROF_SVCS}B",
        color="#1A237E",fs=7.5,ha="left")

    # Stage 2 - total revenue
    rev_b=Y0; rev_t=rev_b+vh(SANKEY_TOTAL_REV)
    ribbon(X[1]+0.10,sub_b,sub_t,X[2]-0.10,rev_b,rev_b+vh(SANKEY_SUB_SUPPORT),"#1565C0",0.32)
    ribbon(X[1]+0.10,pro_b,pro_t,X[2]-0.10,rev_b+vh(SANKEY_SUB_SUPPORT),rev_t,"#283593",0.32)
    nbar(X[2],rev_b,rev_t-rev_b,"#1976D2")
    lbl(X[2],rev_t+0.32,f"Revenue\n${SANKEY_TOTAL_REV}B",color="#0D47A1",fs=11,bold=True)

    # Stage 3 - gross profit + CoR
    cor_b=rev_b; cor_t=cor_b+vh(SANKEY_COR)
    gp_b=cor_t; gp_t=gp_b+vh(SANKEY_GROSS_PROFIT)
    ribbon(X[2]+0.10,rev_b,cor_t,X[3]-0.10,cor_b,cor_t,"#C62828",0.38)
    ribbon(X[2]+0.10,cor_t,rev_t,X[3]-0.10,gp_b,gp_t,"#2E7D32",0.38)
    nbar(X[3],cor_b,cor_t-cor_b,"#C62828"); nbar(X[3],gp_b,gp_t-gp_b,"#388E3C")
    lbl(X[3],gp_t+0.32,
        f"Gross Profit  ${SANKEY_GROSS_PROFIT}B  |  {SANKEY_GROSS_PROFIT/SANKEY_TOTAL_REV*100:.0f}% margin",
        color="#1B5E20",fs=9.5,bold=True)
    lbl(X[3],(cor_b+cor_t)/2,f"Cost of\nRevenue\n(${SANKEY_COR}B)",color="#B71C1C",fs=8)

    # Stage 4 - op profit + opex
    opex_b=gp_b; opex_t=opex_b+vh(SANKEY_OP_EXP)
    op_b=opex_t; op_t=op_b+vh(SANKEY_OP_PROFIT)
    ribbon(X[3]+0.10,gp_b,opex_t,X[4]-0.10,opex_b,opex_t,"#C62828",0.32)
    ribbon(X[3]+0.10,opex_t,gp_t,X[4]-0.10,op_b,op_t,"#2E7D32",0.38)
    ribbon(X[3]+0.10,cor_b,cor_t,X[4]-0.10,cor_b,cor_t,"#C62828",0.18)
    nbar(X[4],opex_b,opex_t-opex_b,"#D32F2F"); nbar(X[4],op_b,op_t-op_b,"#388E3C")
    lbl(X[4]-0.6,(opex_b+opex_t)/2,f"Operating\nExpenses\n(${SANKEY_OP_EXP}B)",
        color="#B71C1C",fs=9,ha="right",bold=True)
    lbl(X[4]-0.6,(op_b+op_t)/2,
        f"Operating Profit\n${SANKEY_OP_PROFIT}B  {SANKEY_OP_PROFIT/SANKEY_TOTAL_REV*100:.0f}%",
        color="#1B5E20",fs=8.5,ha="right",bold=True)

    # Stage 5 - net profit + tax; opex breakdown
    net_b=op_b; net_t=net_b+vh(SANKEY_NET_PROFIT)
    tax_b=net_t; tax_t=tax_b+vh(SANKEY_TAX)
    ribbon(X[4]+0.10,op_b,net_t,X[5]-0.10,net_b,net_t,"#2E7D32",0.48)
    ribbon(X[4]+0.10,net_t,op_t,X[5]-0.10,tax_b,tax_t,"#C62828",0.32)
    nbar(X[5],net_b,net_t-net_b,"#388E3C"); nbar(X[5],tax_b,tax_t-tax_b,"#EF5350")
    lbl(X[5]+0.65,(net_b+net_t)/2,
        f"Net Profit  ${SANKEY_NET_PROFIT}B\n{SANKEY_NET_PROFIT/SANKEY_TOTAL_REV*100:.0f}% net margin",
        color="#1B5E20",fs=9.5,ha="left",bold=True)
    lbl(X[5]+0.65,(tax_b+tax_t)/2,f"Tax  (${SANKEY_TAX}B)",color="#C62828",fs=8,ha="left")
    y_od=opex_b
    for name,v,c in SANKEY_OPEX_DETAIL:
        ribbon(X[4]+0.10,y_od,y_od+vh(v),X[5]-0.10,y_od,y_od+vh(v),c,0.38)
        nbar(X[5],y_od,vh(v),c)
        if vh(v)>0.10:
            lbl(X[5]+0.65,y_od+vh(v)/2,f"{name}  (${v}B)  {v/SANKEY_TOTAL_REV*100:.0f}%",
                color="#B71C1C",fs=7.5,ha="left")
        y_od+=vh(v)
    nbar(X[5],cor_b,cor_t-cor_b,"#FFCDD2")
    lbl(X[5]+0.65,(cor_b+cor_t)/2,f"Cost of Rev.\n(${SANKEY_COR}B)",color="#B71C1C",fs=7.5,ha="left")

    ax.text(W/2,H-0.4,"Income Statement Sankey  -  FY 2024  ($B)",
            ha="center",va="center",fontsize=14,fontweight="bold",color="#0D47A1")

    buf = io.BytesIO()
    fig.savefig(buf,format="png",dpi=150,bbox_inches="tight",facecolor="white")
    plt.close(fig); buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
#  EXCEL WORKBOOK
# ══════════════════════════════════════════════════════════════════════════════

def build_excel():
    try:
        import xlsxwriter
    except ImportError:
        print("  [!] xlsxwriter not found – skipping Excel.  Run: pip install xlsxwriter")
        return

    wb = xlsxwriter.Workbook("McKinsey_Charts.xlsx")

    # ── Shared formats ─────────────────────────────────────────────────────────
    F = {}
    F["title"] = wb.add_format({
        "bold": True, "font_name": "Arial", "font_size": 13,
        "bottom": 2, "bottom_color": DARK_BLUE, "valign": "vcenter",
    })
    F["subtitle"] = wb.add_format({
        "font_name": "Arial", "font_size": 9,
        "font_color": MID_GRAY, "italic": True,
    })
    F["hdr"] = wb.add_format({
        "bold": True, "font_name": "Arial", "font_size": 9,
        "bg_color": DARK_BLUE, "font_color": WHITE,
        "align": "center", "valign": "vcenter", "text_wrap": True, "border": 0,
    })
    F["lbl"] = wb.add_format({"font_name": "Arial", "font_size": 9, "bold": True})
    F["num"] = wb.add_format({
        "font_name": "Arial", "font_size": 9,
        "num_format": "#,##0", "align": "right",
    })
    F["pct"] = wb.add_format({
        "font_name": "Arial", "font_size": 9,
        "num_format": "0%", "align": "right",
    })
    F["src"] = wb.add_format({
        "font_name": "Arial", "font_size": 8,
        "font_color": MID_GRAY, "italic": True,
    })
    F["dpos"] = wb.add_format({
        "font_name": "Arial", "font_size": 9, "bold": True, "font_color": GREEN,
    })
    F["dneg"] = wb.add_format({
        "font_name": "Arial", "font_size": 9, "bold": True, "font_color": RED,
    })

    def sheet(name):
        ws = wb.add_worksheet(name)
        ws.set_zoom(90)
        ws.hide_gridlines(2)
        return ws

    # ── Chart 1: GBS Waterfall ────────────────────────────────────────────────
    def chart1():
        SN = "1 GBS Waterfall"
        ws = sheet(SN)
        ws.set_column("A:A", 0.5)
        ws.set_column("B:B", 20)
        ws.set_column("C:I", 10)

        ws.merge_range("B2:I2", "Global Business Services Scope", F["title"])
        ws.write("B3", "FTE", F["subtitle"])

        # Data table header
        R0, R1 = 4, 5
        n = len(GBS_CATS)
        headers = ["Category", "Spacer", "Reduction", "Accounting",
                   "Customer Service", "Procurement", "Sales Support", "Total"]
        for c, h in enumerate(headers):
            ws.write(R0, 1+c, h, F["hdr"])

        all_rows = list(zip(GBS_CATS, GBS_SPACER, GBS_REDUCTION,
                            GBS_ACCOUNTING, GBS_CUST_SERVICE,
                            GBS_PROCUREMENT, GBS_SALES_SUPP, GBS_TOTALS))
        for r, row in enumerate(all_rows):
            ws.write(R1+r, 1, row[0], F["lbl"])
            for c, v in enumerate(row[1:]):
                ws.write(R1+r, 2+c, v, F["num"])

        chart = wb.add_chart({"type": "column", "subtype": "stacked"})

        # Series 1 – Spacer (invisible, positions reduction bars)
        chart.add_series({
            "name":       "Spacer",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 2, R1+n-1, 2],
            "fill":       {"color": WHITE},
            "border":     {"none": True},
            "data_labels": {"value": False},
        })

        # Series 2 – Reduction bars (gray, floating)
        chart.add_series({
            "name":       "Scope reduction",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 3, R1+n-1, 3],
            "fill":       {"color": LIGHT_GRAY},
            "border":     {"none": True},
            "data_labels": {
                "value": True,
                "font": {"size": 9, "bold": True, "name": "Arial", "color": MID_GRAY},
            },
        })

        # Series 3–6 – Stacked milestone segments
        seg_colors = [BLACK, DARK_BLUE, MID_BLUE, LIGHT_BLUE]
        seg_names  = ["Accounting", "Customer Service", "Procurement", "Sales Support"]
        for i, (seg, color) in enumerate(zip(seg_names, seg_colors)):
            chart.add_series({
                "name":       seg,
                "categories": [SN, R1, 1, R1+n-1, 1],
                "values":     [SN, R1, 4+i, R1+n-1, 4+i],
                "fill":       {"color": color},
                "border":     {"none": True},
                "data_labels": {
                    "value": True,
                    "font": {
                        "size": 9, "bold": True, "name": "Arial",
                        "color": WHITE if color in (BLACK, DARK_BLUE, MID_BLUE) else BLACK,
                    },
                },
            })

        chart.set_title({"none": True})
        chart.set_legend({
            "position": "bottom",
            "font": {"size": 9, "name": "Arial"},
            "delete_series": [0],   # hide "Spacer" from legend
        })
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_x_axis({
            "num_font": {"size": 9, "name": "Arial", "bold": True},
            "line": {"color": MID_GRAY},
            "major_gridlines": {"visible": False},
        })
        chart.set_y_axis({
            "num_font": {"size": 8, "name": "Arial"},
            "line": {"none": True},
            "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}},
            "min": 0, "max": 220,
        })
        chart.set_size({"width": 780, "height": 440})
        ws.insert_chart("B15", chart)
        ws.write(R1+n+22, 1, "Source: Internal HR data", F["src"])

    # ── Chart 2: Diverging Bar ─────────────────────────────────────────────────
    def chart2():
        SN = "2 Diverging Bar"
        ws = sheet(SN)
        ws.set_column("A:A", 0.5)
        ws.set_column("B:B", 34)
        ws.set_column("C:H", 10)

        ws.merge_range("B2:H2", "Purchasing Manager's Sentiment", F["title"])
        ws.write("B3", "% that agree with the following statements", F["subtitle"])

        R0, R1 = 4, 5
        for c, h in enumerate(["Statement","Strongly\nDisagree","Disagree","Agree","Strongly\nAgree","Δ pts"]):
            ws.write(R0, 1+c, h, F["hdr"])

        for r, (stmt, vals) in enumerate(zip(SURVEY_STATEMENTS, SURVEY_DATA)):
            ws.write(R1+r, 1, stmt, F["lbl"])
            for c, v in enumerate(vals[:4]):
                ws.write(R1+r, 2+c, v/100, F["pct"])
            delta = vals[4]
            ws.write(R1+r, 6, f"{'▲' if delta >= 0 else '▼'} {abs(delta)}",
                     F["dpos"] if delta >= 0 else F["dneg"])

        n = len(SURVEY_STATEMENTS)
        chart = wb.add_chart({"type": "bar", "subtype": "stacked"})
        for name, col, color in [
            ("Strongly Disagree", 2, "#C00000"),
            ("Disagree",          3, "#FF7C80"),
            ("Agree",             4, LIGHT_BLUE),
            ("Strongly Agree",    5, DARK_BLUE),
        ]:
            chart.add_series({
                "name":       name,
                "categories": [SN, R1, 1, R1+n-1, 1],
                "values":     [SN, R1, col, R1+n-1, col],
                "fill":       {"color": color},
                "border":     {"none": True},
                "data_labels": {
                    "value": True,
                    "num_format": "0%;0%",
                    "font": {
                        "size": 8, "name": "Arial",
                        "color": WHITE if color in ("#C00000", DARK_BLUE) else BLACK,
                    },
                },
            })
        chart.set_title({"none": True})
        chart.set_legend({"position": "bottom", "font": {"size": 9, "name": "Arial"}})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_x_axis({
            "num_font": {"size": 8, "name": "Arial"},
            "num_format": "0%;0%",
            "line": {"none": True},
            "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}},
            "crossing": 0,
        })
        chart.set_y_axis({
            "num_font": {"size": 8, "name": "Arial"},
            "line": {"none": True},
            "major_gridlines": {"visible": False},
            "reverse": True,
        })
        chart.set_size({"width": 720, "height": 380})
        ws.insert_chart("B13", chart)
        ws.write(R1+n+15, 1, "Source: Annual Procurement Survey 2024  |  n = 245", F["src"])

    # ── Chart 3: Scatter + Line ────────────────────────────────────────────────
    def chart3():
        SN = "3 Scatter Line"
        ws = sheet(SN)
        ws.set_column("A:A", 0.5)
        ws.set_column("B:K", 12)

        ws.merge_range("B2:I2", "Gross Margin Analysis by Promotion Type", F["title"])
        ws.write("B3", "GM cents per unit vs. units sold (indexed, base = 100)", F["subtitle"])

        R0, R1 = 4, 5
        for c, h in enumerate(["30% Off – Units","30% Off – GM","Buy2Get1 – Units","Buy2Get1 – GM","Target Units","Target GM"]):
            ws.write(R0, 1+c, h, F["hdr"])

        for r, ((u1,g1),(u2,g2)) in enumerate(zip(PROMO_30, BUY_2_GET_1)):
            ws.write(R1+r, 1, u1, F["num"])
            ws.write(R1+r, 2, g1, F["num"])
            ws.write(R1+r, 3, u2, F["num"])
            ws.write(R1+r, 4, g2, F["num"])
        # two-point reference line
        ws.write(R1,   5,  80,        F["num"]); ws.write(R1,   6, GM_TARGET, F["num"])
        ws.write(R1+1, 5, 1000,       F["num"]); ws.write(R1+1, 6, GM_TARGET, F["num"])

        n = len(PROMO_30)
        chart = wb.add_chart({"type": "scatter", "subtype": "straight_with_markers"})
        chart.add_series({
            "name":       "30% Off",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 2, R1+n-1, 2],
            "marker":     {"type": "circle",  "size": 7,
                           "fill": {"color": DARK_BLUE}, "border": {"none": True}},
            "line":       {"color": DARK_BLUE, "width": 2.25},
        })
        chart.add_series({
            "name":       "Buy 2 Get 1 Free",
            "categories": [SN, R1, 3, R1+n-1, 3],
            "values":     [SN, R1, 4, R1+n-1, 4],
            "marker":     {"type": "diamond", "size": 7,
                           "fill": {"color": MID_BLUE},  "border": {"none": True}},
            "line":       {"color": MID_BLUE,  "width": 2.25},
        })
        chart.add_series({
            "name":       "Target GM",
            "categories": [SN, R1, 5, R1+1, 5],
            "values":     [SN, R1, 6, R1+1, 6],
            "marker":     {"type": "none"},
            "line":       {"color": RED, "width": 1.5, "dash_type": "dash"},
        })
        chart.set_title({"none": True})
        chart.set_legend({"position": "bottom", "font": {"size": 9, "name": "Arial"}})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_x_axis({
            "name": "Units Sold (indexed)",
            "name_font": {"size": 9, "name": "Arial", "bold": True},
            "num_font": {"size": 8, "name": "Arial"},
            "line": {"color": LIGHT_GRAY},
            "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}},
        })
        chart.set_y_axis({
            "name": "GM (cents per unit)",
            "name_font": {"size": 9, "name": "Arial", "bold": True},
            "num_font": {"size": 8, "name": "Arial"},
            "line": {"none": True},
            "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}},
            "min": 0,
        })
        chart.set_size({"width": 640, "height": 420})
        ws.insert_chart("B13", chart)
        ws.write(R1+n+15, 1, "Source: Trade & Promotions dataset Q1–Q4 2024", F["src"])

    # ── Chart 4: Doughnut ──────────────────────────────────────────────────────
    def chart4():
        SN = "4 Doughnut"
        ws = sheet(SN)
        ws.set_column("A:A", 0.5)
        ws.set_column("B:O", 10)

        ws.merge_range("B2:M2", "Online Revenue by Category", F["title"])
        ws.write("B3", "Share of total revenue, %", F["subtitle"])

        R0, R1 = 4, 5
        start_col = 1   # column B (0-indexed)

        for year, vals in REV_YEARS.items():
            ws.merge_range(R0, start_col, R0, start_col+1, year, F["hdr"])
            for r, (cat, v) in enumerate(zip(REV_CATS, vals)):
                ws.write(R1+r, start_col,   cat, F["lbl"])
                ws.write(R1+r, start_col+1, v,   F["pct"])

            chart = wb.add_chart({"type": "doughnut"})
            chart.add_series({
                "name":       year,
                "categories": [SN, R1, start_col,   R1+len(REV_CATS)-1, start_col],
                "values":     [SN, R1, start_col+1, R1+len(REV_CATS)-1, start_col+1],
                "points":     [{"fill": {"color": c}, "border": {"none": True}}
                               for c in REV_COLORS],
                "data_labels": {
                    "percentage": True,
                    "font": {"size": 8, "name": "Arial"},
                },
            })
            chart.set_title({
                "name":      year,
                "name_font": {"size": 14, "bold": True, "name": "Arial", "color": DARK_BLUE},
            })
            chart.set_legend({"none": True})
            chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
            chart.set_plotarea({"border": {"none": True}})
            chart.set_hole_size(45)
            chart.set_size({"width": 270, "height": 280})
            ws.insert_chart(12, start_col, chart)
            start_col += 4

        # Colour legend
        leg_row = 27
        ws.write(leg_row, 1, "Category legend:", F["lbl"])
        for i, (cat, color) in enumerate(zip(REV_CATS, REV_COLORS)):
            cfmt = wb.add_format({
                "font_name": "Arial", "font_size": 8,
                "bg_color": color,
                "font_color": WHITE if color not in (LIGHT_GRAY, PALE_BLUE, LIGHT_BLUE) else BLACK,
                "align": "center", "border": 0,
            })
            ws.write(leg_row+1, 1+i, cat, cfmt)
        ws.write(leg_row+3, 1, "Source: E-commerce Analytics Platform 2024", F["src"])

    # ── Chart 5: Range Column + Line ──────────────────────────────────────────
    def chart5():
        SN = "5 Range Column"
        ws = sheet(SN)
        ws.set_column("A:A", 0.5)
        ws.set_column("B:J", 10)

        ws.merge_range("B2:J2", "Daily Shipment Volume – Saturdays", F["title"])
        ws.write("B3", "Units (thousands)", F["subtitle"])

        R0, R1 = 4, 5
        for c, h in enumerate(["Week","Min","Range","3yr Avg","Actual","Target"]):
            ws.write(R0, 1+c, h, F["hdr"])

        rng = [mx - mn for mx, mn in zip(SHIP_MAX, SHIP_MIN)]
        for r in range(len(SHIP_WEEKS)):
            ws.write(R1+r, 1, SHIP_WEEKS[r],  F["lbl"])
            ws.write(R1+r, 2, SHIP_MIN[r],    F["num"])
            ws.write(R1+r, 3, rng[r],         F["num"])
            ws.write(R1+r, 4, SHIP_AVG[r],    F["num"])
            ws.write(R1+r, 5, SHIP_ACTUAL[r], F["num"])
            ws.write(R1+r, 6, SHIP_TARGET[r], F["num"])

        n = len(SHIP_WEEKS)

        # Base column chart (range band via stacked)
        col_chart = wb.add_chart({"type": "column", "subtype": "stacked"})
        col_chart.add_series({
            "name":       "Min (hidden)",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 2, R1+n-1, 2],
            "fill":       {"color": WHITE},
            "border":     {"none": True},
        })
        col_chart.add_series({
            "name":       "3-Year Range",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 3, R1+n-1, 3],
            "fill":       {"color": PALE_BLUE},
            "border":     {"none": True},
        })

        # Line chart overlay
        line_chart = wb.add_chart({"type": "line"})
        line_chart.add_series({
            "name":       "3-Year Avg",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 4, R1+n-1, 4],
            "line":       {"color": MID_GRAY, "width": 1.5, "dash_type": "dash"},
            "marker":     {"type": "none"},
        })
        line_chart.add_series({
            "name":       "Actual",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 5, R1+n-1, 5],
            "line":       {"color": DARK_BLUE, "width": 2.5},
            "marker":     {"type": "circle", "size": 7,
                           "fill": {"color": DARK_BLUE}, "border": {"none": True}},
            "data_labels": {
                "value": True, "position": "above",
                "font": {"size": 8, "bold": True, "name": "Arial", "color": DARK_BLUE},
            },
        })
        line_chart.add_series({
            "name":       "Target",
            "categories": [SN, R1, 1, R1+n-1, 1],
            "values":     [SN, R1, 6, R1+n-1, 6],
            "line":       {"color": RED, "width": 1.5, "dash_type": "dash"},
            "marker":     {"type": "none"},
        })

        col_chart.combine(line_chart)
        col_chart.set_title({"none": True})
        col_chart.set_legend({"position": "bottom", "font": {"size": 9, "name": "Arial"}})
        col_chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        col_chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        col_chart.set_x_axis({
            "num_font": {"size": 8, "name": "Arial"},
            "line": {"none": True},
            "major_gridlines": {"visible": False},
        })
        col_chart.set_y_axis({
            "num_font": {"size": 8, "name": "Arial"},
            "line": {"none": True},
            "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}},
            "min": 0,
        })
        col_chart.set_size({"width": 720, "height": 400})
        ws.insert_chart("B13", col_chart)
        ws.write(R1+n+18, 1, "Source: Logistics Operations Data 2024", F["src"])

    # ── Chart 6: Mekko / Marimekko ────────────────────────────────────────────
    def chart6():
        SN = "6 Mekko"
        ws = sheet(SN)
        ws.set_column("A:A", 0.5)
        ws.set_column("B:B", 18)
        ws.set_column("C:F", 14)

        ws.merge_range("B2:F2", "Market Landscape by Region & Product Line", F["title"])
        ws.write("B3", "Column width = market size ($M)  |  Height = share within region (%)", F["subtitle"])

        # Data table
        R0, R1 = 4, 5
        ws.write(R0, 1, "Product Line", F["hdr"])
        for c, seg in enumerate(MEKKO_SEGMENTS):
            ws.write(R0, 2+c, seg, F["hdr"])

        for r, (row_lbl, row_vals) in enumerate(zip(MEKKO_ROWS, MEKKO_DATA)):
            ws.write(R1+r, 1, row_lbl, F["lbl"])
            for c, v in enumerate(row_vals):
                ws.write(R1+r, 2+c, v/100, F["pct"])

        # Market size row
        ws.write(R1+len(MEKKO_ROWS)+1, 1, "Market size ($M)", F["lbl"])
        for c, sz in enumerate(MEKKO_SEG_SIZES):
            ws.write(R1+len(MEKKO_ROWS)+1, 2+c, sz, F["num"])

        # 100% stacked column chart (Excel approximation of Mekko – widths equal)
        chart = wb.add_chart({"type": "column", "subtype": "percent_stacked"})
        n_rows = len(MEKKO_ROWS)
        for i, (row_lbl, color) in enumerate(zip(MEKKO_ROWS, MEKKO_ROW_COLORS)):
            chart.add_series({
                "name":       [SN, R1+i, 1],
                "categories": [SN, R0,   2, R0,   2+len(MEKKO_SEGMENTS)-1],
                "values":     [SN, R1+i, 2, R1+i, 2+len(MEKKO_SEGMENTS)-1],
                "fill":       {"color": color},
                "border":     {"none": True},
                "data_labels": {
                    "percentage": True,
                    "font": {
                        "size": 9, "bold": True, "name": "Arial",
                        "color": WHITE if color in (DARK_BLUE, MID_BLUE, BLACK) else BLACK,
                    },
                },
            })

        # Annotation: market sizes above chart
        note_fmt = wb.add_format({"font_name": "Arial", "font_size": 8,
                                   "font_color": DARK_BLUE, "bold": True, "align": "center"})
        for c, sz in enumerate(MEKKO_SEG_SIZES):
            ws.write(R0-1, 2+c, f"${sz}M", note_fmt)

        chart.set_title({"none": True})
        chart.set_legend({"position": "bottom", "font": {"size": 9, "name": "Arial"}})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_x_axis({
            "num_font": {"size": 9, "name": "Arial", "bold": True},
            "line": {"none": True},
            "major_gridlines": {"visible": False},
        })
        chart.set_y_axis({
            "num_font": {"size": 8, "name": "Arial"},
            "num_format": "0%",
            "line": {"none": True},
            "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}},
        })
        chart.set_size({"width": 640, "height": 420})
        ws.insert_chart("B13", chart)
        ws.write(R1+n_rows+10, 1, "Note: Column widths equal in Excel; see PowerPoint for proportional Mekko", F["src"])
        ws.write(R1+n_rows+11, 1, "Source: Market Intelligence Report 2024", F["src"])

    # ─────────────────────────────────────────────────────────────────────────
    # Sheet 7: 2×2 Priority Matrix
    # ─────────────────────────────────────────────────────────────────────────
    def chart7():
        SN = "7_Priority_Matrix"
        ws = wb.add_worksheet(SN)
        ws.set_tab_color(MID_BLUE)
        ws.set_column(0, 0, 24)
        ws.set_column(1, 2, 12)
        ws.write(0, 0, "2×2 Priority Matrix", F["title"])
        ws.write(1, 0, "Initiative", F["hdr"])
        ws.write(1, 1, "Impl. Ease (0-10)", F["hdr"])
        ws.write(1, 2, "Business Impact (0-10)", F["hdr"])
        for r, (lbl, ease, impact) in enumerate(MATRIX_ITEMS):
            ws.write(2+r, 0, lbl.replace("\n", " "), F["lbl"])
            ws.write(2+r, 1, ease, F["num"])
            ws.write(2+r, 2, impact, F["num"])
        chart = wb.add_chart({"type": "scatter"})
        chart.add_series({
            "name":       "Initiatives",
            "categories": [SN, 2, 1, 2+len(MATRIX_ITEMS)-1, 1],
            "values":     [SN, 2, 2, 2+len(MATRIX_ITEMS)-1, 2],
            "marker": {"type": "circle", "size": 8,
                       "fill": {"color": DARK_BLUE}, "border": {"color": WHITE}},
        })
        # Quadrant reference lines
        for dummy_y, ref_name in [(5, "Ease=5"), (5, "Impact=5")]:
            pass  # reference lines via series
        chart.add_series({
            "name": "_ref_v",
            "categories": [SN, 2, 1, 2, 1],
            "values":     [SN, 2, 1, 2, 1],
            "line": {"color": MID_GRAY, "dash_type": "dash", "width": 1},
            "marker": {"type": "none"},
        })
        chart.set_title({"name": "2×2 Initiative Priority Matrix"})
        chart.set_x_axis({"name": "Implementation Ease →", "min": 0, "max": 10,
                          "crossing": 5,
                          "num_font": {"size": 9, "name": "Arial"},
                          "name_font": {"size": 10, "name": "Arial", "bold": True}})
        chart.set_y_axis({"name": "Business Impact →", "min": 0, "max": 10,
                          "crossing": 5,
                          "num_font": {"size": 9, "name": "Arial"},
                          "name_font": {"size": 10, "name": "Arial", "bold": True}})
        chart.set_legend({"none": True})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"color": LIGHT_GRAY}, "fill": {"color": WHITE}})
        chart.set_size({"width": 520, "height": 420})
        ws.insert_chart("E2", chart)
        ws.write(2+len(MATRIX_ITEMS)+1, 0, "Source: Strategic Planning Workshop 2024", F["src"])

    # ─────────────────────────────────────────────────────────────────────────
    # Sheet 8: EBITDA Bridge
    # ─────────────────────────────────────────────────────────────────────────
    def chart8():
        SN = "8_EBITDA_Bridge"
        ws = wb.add_worksheet(SN)
        ws.set_tab_color(DARK_BLUE)
        ws.set_column(0, 0, 22)
        ws.set_column(1, 3, 10)
        ws.write(0, 0, "EBITDA Bridge", F["title"])
        ws.write(1, 0, "Item", F["hdr"])
        ws.write(1, 1, "Value ($M)", F["hdr"])
        ws.write(1, 2, "Spacer", F["hdr"])
        ws.write(1, 3, "Bar", F["hdr"])

        running = 0
        spacers, bars, colors_list, labels_list = [], [], [], []
        for lbl_txt, val, kind in EBITDA_ITEMS:
            if kind == "start":
                spacers.append(0);      bars.append(val);  running = val
                colors_list.append(DARK_BLUE)
            elif kind == "total":
                spacers.append(0);      bars.append(running)
                colors_list.append(DARK_BLUE)
                running = running  # stays
            elif kind == "up":
                spacers.append(running); bars.append(val);  running += val
                colors_list.append("#70AD47")
            else:  # down
                running += val
                spacers.append(running); bars.append(-val)
                colors_list.append("#C00000")
            labels_list.append(lbl_txt)

        for r, (lbl_txt, sp, ba, val_orig) in enumerate(zip(labels_list, spacers, bars, [v for _,v,_ in EBITDA_ITEMS])):
            ws.write(2+r, 0, lbl_txt, F["lbl"])
            ws.write(2+r, 1, val_orig, F["num"])
            ws.write(2+r, 2, sp, F["num"])
            ws.write(2+r, 3, ba, F["num"])

        chart = wb.add_chart({"type": "bar", "subtype": "stacked"})
        chart.add_series({
            "name": "Spacer",
            "categories": [SN, 2, 0, 2+len(EBITDA_ITEMS)-1, 0],
            "values":     [SN, 2, 2, 2+len(EBITDA_ITEMS)-1, 2],
            "fill": {"none": True}, "border": {"none": True},
        })
        chart.add_series({
            "name": "Value",
            "categories": [SN, 2, 0, 2+len(EBITDA_ITEMS)-1, 0],
            "values":     [SN, 2, 3, 2+len(EBITDA_ITEMS)-1, 3],
            "fill": {"color": DARK_BLUE}, "border": {"none": True},
            "data_labels": {"value": True, "font": {"size": 9, "bold": True, "name": "Arial", "color": WHITE}},
        })
        chart.set_title({"name": "EBITDA Bridge  ($M)"})
        chart.set_x_axis({"num_font": {"size": 9, "name": "Arial"}, "line": {"none": True},
                          "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}}})
        chart.set_y_axis({"num_font": {"size": 9, "name": "Arial"}, "line": {"none": True},
                          "major_gridlines": {"visible": False}})
        chart.set_legend({"none": True})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_size({"width": 560, "height": 400})
        ws.insert_chart("F2", chart)
        ws.write(2+len(EBITDA_ITEMS)+1, 0, "Source: Management Accounts FY2024", F["src"])

    # ─────────────────────────────────────────────────────────────────────────
    # Sheet 9: Gantt / Roadmap
    # ─────────────────────────────────────────────────────────────────────────
    def chart9():
        SN = "9_Gantt_Roadmap"
        ws = wb.add_worksheet(SN)
        ws.set_tab_color(MID_BLUE)
        ws.set_column(0, 0, 22)
        ws.set_column(1, 1, 26)
        ws.set_column(2, 3, 10)
        ws.write(0, 0, "Transformation Roadmap", F["title"])
        ws.write(1, 0, "Workstream", F["hdr"])
        ws.write(1, 1, "Task", F["hdr"])
        ws.write(1, 2, "Start (week)", F["hdr"])
        ws.write(1, 3, "Duration (weeks)", F["hdr"])
        for r, (ws_name, task, start, dur, color) in enumerate(GANTT_DATA):
            ws.write(2+r, 0, ws_name, F["lbl"])
            ws.write(2+r, 1, task, F["lbl"])
            ws.write(2+r, 2, start, F["num"])
            ws.write(2+r, 3, dur, F["num"])
        chart = wb.add_chart({"type": "bar", "subtype": "stacked"})
        chart.add_series({
            "name": "Start",
            "categories": [SN, 2, 1, 2+len(GANTT_DATA)-1, 1],
            "values":     [SN, 2, 2, 2+len(GANTT_DATA)-1, 2],
            "fill": {"none": True}, "border": {"none": True},
        })
        chart.add_series({
            "name": "Duration",
            "categories": [SN, 2, 1, 2+len(GANTT_DATA)-1, 1],
            "values":     [SN, 2, 3, 2+len(GANTT_DATA)-1, 3],
            "fill": {"color": DARK_BLUE}, "border": {"none": True},
            "data_labels": {"value": False},
        })
        chart.set_title({"name": "Transformation Roadmap  (weeks)"})
        chart.set_x_axis({"name": "Week", "min": 0, "max": GANTT_TOTAL_WEEKS,
                          "num_font": {"size": 9, "name": "Arial"},
                          "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}}})
        chart.set_y_axis({"num_font": {"size": 9, "name": "Arial"}, "line": {"none": True},
                          "major_gridlines": {"visible": False}})
        chart.set_legend({"none": True})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_size({"width": 620, "height": 400})
        ws.insert_chart("F2", chart)
        ws.write(2+len(GANTT_DATA)+1, 0, "Source: Programme Management Office", F["src"])

    # ─────────────────────────────────────────────────────────────────────────
    # Sheet 10: Harvey Balls
    # ─────────────────────────────────────────────────────────────────────────
    def chart10():
        SYMBOLS = ["○", "◔", "◑", "◕", "●"]
        SN = "10_Harvey_Balls"
        ws = wb.add_worksheet(SN)
        ws.set_tab_color(LIGHT_BLUE)
        ws.set_column(0, 0, 26)
        for c in range(1, len(HARVEY_COMPANIES)+1):
            ws.set_column(c, c, 14)
        ws.set_row(0, 22)
        ws.write(0, 0, "Capability Benchmarking", F["title"])
        ws.write(1, 0, "Capability", F["hdr"])
        big_fmt = wb.add_format({
            "font_name": "Arial", "font_size": 18,
            "align": "center", "valign": "vcenter",
            "font_color": DARK_BLUE, "bold": True, "border": 1, "border_color": LIGHT_GRAY,
        })
        hdr_fmt = wb.add_format({
            "font_name": "Arial", "font_size": 10, "bold": True,
            "align": "center", "font_color": WHITE, "bg_color": DARK_BLUE,
            "border": 1, "border_color": WHITE,
        })
        crit_fmt = wb.add_format({
            "font_name": "Arial", "font_size": 10,
            "align": "left", "valign": "vcenter",
            "border": 1, "border_color": LIGHT_GRAY,
        })
        for c, co in enumerate(HARVEY_COMPANIES):
            ws.write(1, 1+c, co, hdr_fmt)
        for r, (crit, row_scores) in enumerate(zip(HARVEY_CRITERIA, HARVEY_SCORES)):
            ws.set_row(2+r, 26)
            ws.write(2+r, 0, crit, crit_fmt)
            for c, score in enumerate(row_scores):
                ws.write(2+r, 1+c, SYMBOLS[score], big_fmt)
        # Legend
        legend_row = 2 + len(HARVEY_CRITERIA) + 1
        ws.write(legend_row, 0, "Legend:", F["lbl"])
        for i, (s, desc) in enumerate(zip(SYMBOLS, ["No capability","Basic","Moderate","Strong","Leading"])):
            ws.write(legend_row, 1+i, f"{s} {desc}", wb.add_format(
                {"font_name": "Arial", "font_size": 10, "align": "center",
                 "font_color": DARK_BLUE}))
        ws.write(legend_row+1, 0, "Source: Internal capability assessment", F["src"])

    # ─────────────────────────────────────────────────────────────────────────
    # Sheet 11: Tornado / Sensitivity
    # ─────────────────────────────────────────────────────────────────────────
    def chart11():
        SN = "11_Tornado"
        ws = wb.add_worksheet(SN)
        ws.set_tab_color(DARK_BLUE)
        ws.set_column(0, 0, 26)
        ws.set_column(1, 2, 12)
        ws.write(0, 0, "Sensitivity Analysis – NPV Impact ($M)", F["title"])
        ws.write(1, 0, "Driver", F["hdr"])
        ws.write(1, 1, "Downside ($M)", F["hdr"])
        ws.write(1, 2, "Upside ($M)", F["hdr"])
        # Sort by total range descending
        sorted_drivers = sorted(TORNADO_DRIVERS, key=lambda x: abs(x[1])+abs(x[2]), reverse=True)
        for r, (drv, low, high) in enumerate(sorted_drivers):
            ws.write(2+r, 0, drv, F["lbl"])
            ws.write(2+r, 1, low, F["num"])
            ws.write(2+r, 2, high, F["num"])
        chart = wb.add_chart({"type": "bar"})
        chart.add_series({
            "name": "Downside",
            "categories": [SN, 2, 0, 2+len(sorted_drivers)-1, 0],
            "values":     [SN, 2, 1, 2+len(sorted_drivers)-1, 1],
            "fill": {"color": RED}, "border": {"none": True},
            "data_labels": {"value": True, "font": {"size": 8, "bold": True, "name": "Arial", "color": WHITE}},
        })
        chart.add_series({
            "name": "Upside",
            "categories": [SN, 2, 0, 2+len(sorted_drivers)-1, 0],
            "values":     [SN, 2, 2, 2+len(sorted_drivers)-1, 2],
            "fill": {"color": GREEN}, "border": {"none": True},
            "data_labels": {"value": True, "font": {"size": 8, "bold": True, "name": "Arial", "color": WHITE}},
        })
        chart.set_title({"name": "NPV Sensitivity  ($M)  –  Base Case = $217M"})
        chart.set_x_axis({"num_font": {"size": 9, "name": "Arial"}, "crossing": 0,
                          "major_gridlines": {"visible": True, "line": {"color": LIGHT_GRAY, "width": 0.5}}})
        chart.set_y_axis({"num_font": {"size": 9, "name": "Arial"}, "line": {"none": True},
                          "major_gridlines": {"visible": False}})
        chart.set_legend({"position": "bottom", "font": {"size": 9, "name": "Arial"}})
        chart.set_chartarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_plotarea({"border": {"none": True}, "fill": {"color": WHITE}})
        chart.set_size({"width": 560, "height": 420})
        ws.insert_chart("D2", chart)
        ws.write(2+len(sorted_drivers)+1, 0, "Source: Financial model – sensitivity run", F["src"])

    # ─────────────────────────────────────────────────────────────────────────
    # Sheet 12: Sankey (matplotlib image)
    # ─────────────────────────────────────────────────────────────────────────
    def chart12():
        SN = "12_Sankey"
        ws = wb.add_worksheet(SN)
        ws.set_tab_color(MID_BLUE)
        ws.write(0, 0, "Income Statement Sankey", F["title"])
        png = _make_sankey_png()
        if png:
            import io as _io
            ws.insert_image("A3", "sankey.png", {"image_data": png, "x_scale": 0.85, "y_scale": 0.85})
        else:
            ws.write(2, 0, "[Install matplotlib to see Sankey: pip install matplotlib]", F["src"])
        ws.write(60, 0, "Source: Company financials FY2024", F["src"])

    chart1(); chart2(); chart3(); chart4(); chart5(); chart6(); chart7(); chart8(); chart9(); chart10(); chart11(); chart12()
    wb.close()
    print("  ✓  McKinsey_Charts.xlsx  created")


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

def _hex(h):
    """Convert '#RRGGBB' → pptx RGBColor."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def build_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.chart.data import ChartData, BubbleChartData, XyChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
        import pptx.oxml.ns as ns
        from lxml import etree
    except ImportError:
        print("  [!] python-pptx not found – skipping PowerPoint.  Run: pip install python-pptx")
        return

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

    W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ── Helper: blank slide layout ─────────────────────────────────────────────
    def blank_slide():
        layout = prs.slide_layouts[6]   # blank
        return prs.slides.add_slide(layout)

    # ── Helper: add title block (McKinsey style) ───────────────────────────────
    def add_title(slide, title_text, subtitle_text=""):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        # Blue rule bar at top
        bar = slide.shapes.add_shape(
            1,   # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.4), Inches(0.3),
            Inches(12.5), Pt(3),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _hex(DARK_BLUE)
        bar.line.fill.background()

        # Title text
        txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(10), Inches(0.6))
        tf  = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold      = True
        p.font.size      = Pt(20)
        p.font.name      = "Arial"
        p.font.color.rgb = _hex(DARK_BLUE)

        # Subtitle
        if subtitle_text:
            txb2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(10), Inches(0.35))
            tf2  = txb2.text_frame
            p2   = tf2.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size      = Pt(11)
            p2.font.name      = "Arial"
            p2.font.color.rgb = _hex(MID_GRAY)
            p2.font.italic    = True

    # ── Helper: source line ────────────────────────────────────────────────────
    def add_source(slide, text):
        txb = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12), Inches(0.3))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        p.text = text
        p.font.size      = Pt(8)
        p.font.name      = "Arial"
        p.font.color.rgb = _hex(MID_GRAY)
        p.font.italic    = True

    # ── Helper: set series colour ──────────────────────────────────────────────
    def series_color(series, hex_color):
        from pptx.dml.color import RGBColor
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _hex(hex_color)
        series.format.line.color.rgb      = _hex(hex_color)

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 1: GBS Waterfall  (hand-drawn: rectangles + dotted connectors)
    # ─────────────────────────────────────────────────────────────────────────
    def slide1():
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _etree

        slide = blank_slide()
        add_title(slide, "Global Business Services Scope", "FTE")

        # ── Chart geometry ────────────────────────────────────────────────
        CL = Inches(2.0)    # left edge (room for row labels)
        CT = Inches(1.5)    # top edge
        CH = Inches(4.9)    # chart height
        Y_MAX = 220

        def vy(v):
            """Convert FTE value → Y coordinate on slide (integer EMUs)."""
            return int(CT + CH * (1.0 - v / Y_MAX))

        # ── Column layout ─────────────────────────────────────────────────
        BW  = Inches(1.4)   # milestone bar width
        RW  = Inches(0.9)   # reduction bar width
        G1  = Inches(0.30)  # gap: milestone → reduction
        G2  = Inches(0.18)  # gap: reduction → reduction
        G3  = Inches(0.30)  # gap: reduction → milestone
        G4  = Inches(0.30)  # gap: milestone → milestone

        bm_x   = CL
        proc_x = bm_x   + BW + G1
        str_x  = proc_x + RW + G2
        sys_x  = str_x  + RW + G2
        so_x   = sys_x  + RW + G3
        tr_x   = so_x   + BW + G4
        tb_x   = tr_x   + BW + G4

        SEG_DATA   = [GBS_ACCOUNTING, GBS_CUST_SERVICE,
                      GBS_PROCUREMENT, GBS_SALES_SUPP]
        SEG_COLORS = [BLACK, DARK_BLUE, MID_BLUE, LIGHT_BLUE]
        SEG_LABELS = ["Accounting", "Customer Service",
                      "Procurement", "Sales Support"]

        def add_label(bx, bw, bar_top, bar_h, text, fg):
            if bar_h < Inches(0.22):
                return
            txb = slide.shapes.add_textbox(
                bx, int(bar_top + bar_h * 0.08),
                bw, int(bar_h * 0.84))
            p = txb.text_frame.paragraphs[0]
            p.text = str(text)
            p.font.size      = Pt(11)
            p.font.bold      = True
            p.font.name      = "Arial"
            p.font.color.rgb = _hex(fg)
            p.alignment      = PP_ALIGN.CENTER

        # ── Draw milestone stacked columns ────────────────────────────────
        for col_idx, (bx, di) in enumerate([
            (bm_x, 0), (so_x, 4), (tr_x, 5), (tb_x, 6)
        ]):
            cum = 0
            for seg_vals, color in zip(SEG_DATA, SEG_COLORS):
                v = seg_vals[di]
                if v <= 0:
                    cum += v
                    continue
                bar_top = vy(cum + v)
                bar_h   = vy(cum) - bar_top
                rect = slide.shapes.add_shape(1, bx, bar_top, BW, bar_h)
                rect.fill.solid()
                rect.fill.fore_color.rgb = _hex(color)
                rect.line.fill.background()
                fg = WHITE if color in (BLACK, DARK_BLUE, MID_BLUE) else BLACK
                add_label(bx, BW, bar_top, bar_h, v, fg)
                cum += v

            # Total annotation above column
            total = GBS_TOTALS[di]
            ttxb = slide.shapes.add_textbox(
                bx, vy(total) - Inches(0.32), BW, Inches(0.30))
            tp = ttxb.text_frame.paragraphs[0]
            tp.text = str(total)
            tp.font.size      = Pt(12)
            tp.font.bold      = True
            tp.font.name      = "Arial"
            tp.font.color.rgb = _hex(DARK_GRAY)
            tp.alignment      = PP_ALIGN.CENTER

        # ── Draw floating reduction bars ──────────────────────────────────
        for bx, di in [(proc_x, 1), (str_x, 2), (sys_x, 3)]:
            sp  = GBS_SPACER[di]
            red = GBS_REDUCTION[di]
            bar_top = vy(sp + red)
            bar_h   = vy(sp) - bar_top
            rect = slide.shapes.add_shape(1, bx, bar_top, RW, bar_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = _hex(LIGHT_GRAY)
            rect.line.fill.background()
            add_label(bx, RW, bar_top, bar_h, red, MID_GRAY)

        # ── X-axis line ───────────────────────────────────────────────────
        ax_end = tb_x + BW + Inches(0.15)
        ax = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, CL - Inches(0.1), CT + CH,
            ax_end, CT + CH)
        ax.line.color.rgb = _hex(MID_GRAY)
        ax.line.width     = Pt(1)

        # ── X-axis category labels ────────────────────────────────────────
        cat_defs = [
            (bm_x,   BW, "Benchmark\nscope",          False),
            (proc_x, RW, "Process\nstandardization",  False),
            (str_x,  RW, "Strategic\nimportance",      False),
            (sys_x,  RW, "System\nlimitations",        False),
            (so_x,   BW, "Signed-off\nscope",          False),
            (tr_x,   BW, "Transferred\nto date",       False),
            (tb_x,   BW, "To be\ntransferred",         True),
        ]
        for bx, bw, lbl, bold in cat_defs:
            xlbl = slide.shapes.add_textbox(
                bx - Inches(0.08), CT + CH + Inches(0.1),
                bw + Inches(0.16), Inches(0.55))
            xlbl.text_frame.word_wrap = True
            xp = xlbl.text_frame.paragraphs[0]
            xp.text           = lbl
            xp.font.size      = Pt(9)
            xp.font.name      = "Arial"
            xp.font.bold      = bold
            xp.font.color.rgb = _hex(BLACK)
            xp.alignment      = PP_ALIGN.CENTER

        # ── Row labels (left, aligned to Benchmark segments) ─────────────
        cum = 0
        for seg_vals, lbl, color in zip(SEG_DATA, SEG_LABELS, SEG_COLORS):
            v = seg_vals[0]
            if v > 0:
                y_mid = vy(cum + v // 2)
                lt = slide.shapes.add_textbox(
                    CL - Inches(1.5), y_mid - Inches(0.22),
                    Inches(1.45), Inches(0.44))
                lt.text_frame.word_wrap = True
                lp = lt.text_frame.paragraphs[0]
                lp.text           = lbl
                lp.font.size      = Pt(9)
                lp.font.name      = "Arial"
                lp.font.color.rgb = _hex(DARK_GRAY)
                lp.alignment      = PP_ALIGN.RIGHT
            cum += v

        # ── Dotted connector lines ────────────────────────────────────────
        def dotted_line(x1, y1, x2, y2):
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
            conn.line.color.rgb = _hex(MID_GRAY)
            conn.line.width     = Pt(0.75)
            try:
                ln_el = conn.line._ln
                for old in ln_el.findall(_qn('a:prstDash')):
                    ln_el.remove(old)
                pd = _etree.SubElement(ln_el, _qn('a:prstDash'))
                pd.set('val', 'dot')
            except Exception:
                pass

        dotted_line(bm_x   + BW, vy(193), proc_x,      vy(193))  # 193 → Process top
        dotted_line(proc_x + RW, vy(172), str_x,        vy(172))  # 172 → Strategic top
        dotted_line(str_x  + RW, vy(151), sys_x,        vy(151))  # 151 → System top
        dotted_line(sys_x  + RW, vy(129), so_x,         vy(129))  # 129 → Signed-off top
        dotted_line(so_x   + BW, vy(95),  tr_x,         vy(95))   #  95 → Transferred top
        dotted_line(tr_x   + BW, vy(34),  tb_x,         vy(34))   #  34 → To be transferred top

        add_source(slide, "Source: Internal HR data")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 2: Diverging Bar  (Survey Sentiment)
    # ─────────────────────────────────────────────────────────────────────────
    def slide2():
        from pptx.enum.chart import XL_CHART_TYPE
        slide = blank_slide()
        add_title(slide, "Purchasing Manager's Sentiment",
                  "% that agree with the following statements")

        cd = ChartData()
        cd.categories = SURVEY_STATEMENTS
        labels = ["Strongly Disagree", "Disagree", "Agree", "Strongly Agree"]
        series_vals = list(zip(*[[v/100 for v in row[:4]] for row in SURVEY_DATA]))
        for lbl, vals in zip(labels, series_vals):
            cd.add_series(lbl, vals)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(0.4), Inches(1.35),
            Inches(11.5), Inches(5.5),
            cd,
        ).chart

        colors = ["#C00000", "#FF7C80", LIGHT_BLUE, DARK_BLUE]
        for s, color in zip(chart.series, colors):
            series_color(s, color)
            s.data_labels.show_percentage = True
            s.data_labels.font.size = Pt(8)
            s.data_labels.font.name = "Arial"

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = False

        # Delta arrows (text box on right)
        for i, row in enumerate(SURVEY_DATA):
            delta = row[4]
            txb = slide.shapes.add_textbox(
                Inches(11.9), Inches(1.6 + i*0.78), Inches(1.2), Inches(0.5)
            )
            tf = txb.text_frame
            p  = tf.paragraphs[0]
            p.text = f"{'▲' if delta >= 0 else '▼'} {abs(delta)}"
            p.font.size      = Pt(11)
            p.font.bold      = True
            p.font.name      = "Arial"
            p.font.color.rgb = _hex(GREEN if delta >= 0 else RED)

        add_source(slide, "Source: Annual Procurement Survey 2024  |  n = 245")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 3: Scatter + Line  (GM Analysis)
    # ─────────────────────────────────────────────────────────────────────────
    def slide3():
        from pptx.enum.chart import XL_CHART_TYPE
        slide = blank_slide()
        add_title(slide, "Gross Margin Analysis by Promotion Type",
                  "GM cents per unit vs. units sold (indexed, base = 100)")

        cd = XyChartData()
        s1 = cd.add_series("30% Off")
        for u, g in PROMO_30:
            s1.add_data_point(u, g)
        s2 = cd.add_series("Buy 2 Get 1 Free")
        for u, g in BUY_2_GET_1:
            s2.add_data_point(u, g)
        s3 = cd.add_series("Target GM")
        s3.add_data_point(80,   GM_TARGET)
        s3.add_data_point(1000, GM_TARGET)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES,
            Inches(0.8), Inches(1.35),
            Inches(11.5), Inches(5.5),
            cd,
        ).chart

        from pptx.util import Pt
        s_colors = [DARK_BLUE, MID_BLUE, RED]
        for s, color in zip(chart.series, s_colors):
            series_color(s, color)

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = False

        add_source(slide, "Source: Trade & Promotions dataset Q1–Q4 2024")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 4: Doughnut charts  (Revenue Mix)
    # ─────────────────────────────────────────────────────────────────────────
    def slide4():
        from pptx.enum.chart import XL_CHART_TYPE
        slide = blank_slide()
        add_title(slide, "Online Revenue by Category",
                  "Share of total revenue, %")

        positions = [Inches(0.4), Inches(4.55), Inches(8.7)]
        for (year, vals), x_pos in zip(REV_YEARS.items(), positions):
            cd = ChartData()
            cd.categories = REV_CATS
            cd.add_series(year, vals)

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT,
                x_pos, Inches(1.35),
                Inches(4.1), Inches(4.8),
                cd,
            ).chart

            # Colour individual slices
            plot = chart.plots[0]
            for point, color in zip(plot.series[0].points, REV_COLORS):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = _hex(color)

            chart.has_title = True
            chart.chart_title.text_frame.text = year
            chart.chart_title.text_frame.paragraphs[0].font.size  = Pt(18)
            chart.chart_title.text_frame.paragraphs[0].font.bold  = True
            chart.chart_title.text_frame.paragraphs[0].font.name  = "Arial"
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = _hex(DARK_BLUE)
            chart.has_legend = False

        # Shared legend
        for i, (cat, color) in enumerate(zip(REV_CATS, REV_COLORS)):
            txb = slide.shapes.add_textbox(
                Inches(0.4 + i*2.0), Inches(6.6), Inches(1.9), Inches(0.35)
            )
            txb.fill.solid()
            txb.fill.fore_color.rgb = _hex(color)
            tf = txb.text_frame
            p  = tf.paragraphs[0]
            p.text = cat
            p.font.size      = Pt(9)
            p.font.name      = "Arial"
            p.font.color.rgb = _hex(WHITE if color not in (LIGHT_GRAY,PALE_BLUE,LIGHT_BLUE) else BLACK)
            p.alignment      = PP_ALIGN.CENTER

        add_source(slide, "Source: E-commerce Analytics Platform 2024")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 5: Range Column + Line  (Shipment Volume)
    # ─────────────────────────────────────────────────────────────────────────
    def slide5():
        from pptx.enum.chart import XL_CHART_TYPE
        slide = blank_slide()
        add_title(slide, "Daily Shipment Volume – Saturdays",
                  "Units (thousands)")

        rng = [mx - mn for mx, mn in zip(SHIP_MAX, SHIP_MIN)]

        cd = ChartData()
        cd.categories = SHIP_WEEKS
        cd.add_series("Min (hidden)",  SHIP_MIN)
        cd.add_series("3-Year Range",  rng)
        cd.add_series("3-Year Avg",    SHIP_AVG)
        cd.add_series("Actual",        SHIP_ACTUAL)
        cd.add_series("Target",        SHIP_TARGET)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(0.4), Inches(1.35),
            Inches(12.5), Inches(5.5),
            cd,
        ).chart

        # Style series
        hidden = chart.series[0]
        hidden.format.fill.background()
        hidden.format.line.fill.background()

        range_s = chart.series[1]
        range_s.format.fill.solid()
        range_s.format.fill.fore_color.rgb = _hex(PALE_BLUE)
        range_s.format.line.fill.background()

        for idx, color in enumerate([MID_GRAY, DARK_BLUE, RED], start=2):
            series_color(chart.series[idx], color)

        chart.series[3].data_labels.show_value = True
        chart.series[3].data_labels.font.size  = Pt(8)
        chart.series[3].data_labels.font.bold  = True
        chart.series[3].data_labels.font.name  = "Arial"

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.has_title = False
        chart.value_axis.has_gridlines = True

        add_source(slide, "Source: Logistics Operations Data 2024")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 6: Mekko / Marimekko  (proper variable-width columns)
    # ─────────────────────────────────────────────────────────────────────────
    def slide6():
        slide = blank_slide()
        add_title(slide, "Market Landscape by Region & Product Line",
                  "Column width = market size ($M)  |  Height = share within region (%)")

        # Chart drawing area on slide
        CX = Inches(0.9)          # chart left edge
        CY = Inches(1.4)          # chart top edge
        CW = Inches(11.5)         # chart total width
        CH = Inches(5.0)          # chart total height

        total_mkt = sum(MEKKO_SEG_SIZES)

        # Cumulative x offsets
        cum_x = 0
        for j, (seg, sz, seg_col) in enumerate(zip(MEKKO_SEGMENTS, MEKKO_SEG_SIZES, MEKKO_SEG_COLORS)):
            col_w   = CW * sz / total_mkt
            col_x   = CX + CW * cum_x / total_mkt
            cum_pct = 0   # cumulative % from top → rows stacked top to bottom

            # Draw each row cell
            cum_row_pct = 0
            for i, (row_lbl, color) in enumerate(zip(MEKKO_ROWS, MEKKO_ROW_COLORS)):
                pct   = MEKKO_DATA[i][j]
                cell_h = CH * pct / 100
                cell_y = CY + CH * cum_row_pct / 100

                rect = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    col_x, cell_y, col_w - Inches(0.04), cell_h,
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = _hex(color)
                rect.line.fill.background()

                # Label inside cell (only if tall enough)
                if pct >= 8:
                    txb = slide.shapes.add_textbox(
                        col_x + Inches(0.05), cell_y + Inches(0.03),
                        col_w - Inches(0.1),  cell_h - Inches(0.06),
                    )
                    tf = txb.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    p.text = f"{pct}%"
                    p.font.size      = Pt(10)
                    p.font.bold      = True
                    p.font.name      = "Arial"
                    p.font.color.rgb = _hex(
                        WHITE if color in (BLACK, DARK_BLUE, MID_BLUE) else BLACK
                    )

                cum_row_pct += pct

            # Segment label below chart
            lbl_txb = slide.shapes.add_textbox(
                col_x, CY + CH + Inches(0.1), col_w, Inches(0.5),
            )
            lf = lbl_txb.text_frame
            lf.word_wrap = True
            p = lf.paragraphs[0]
            p.text = f"{seg}\n${sz}M"
            p.font.size      = Pt(9)
            p.font.bold      = True
            p.font.name      = "Arial"
            p.font.color.rgb = _hex(DARK_BLUE)
            p.alignment      = PP_ALIGN.CENTER

            cum_x += sz

        # Row legend on right
        for i, (row_lbl, color) in enumerate(zip(MEKKO_ROWS, MEKKO_ROW_COLORS)):
            leg_txb = slide.shapes.add_textbox(
                CX + CW + Inches(0.15), CY + Inches(i * 0.4),
                Inches(1.5), Inches(0.35),
            )
            leg_txb.fill.solid()
            leg_txb.fill.fore_color.rgb = _hex(color)
            lp = leg_txb.text_frame.paragraphs[0]
            lp.text = row_lbl
            lp.font.size      = Pt(9)
            lp.font.name      = "Arial"
            lp.font.color.rgb = _hex(
                WHITE if color in (BLACK, DARK_BLUE, MID_BLUE) else BLACK
            )
            lp.alignment = PP_ALIGN.CENTER

        add_source(slide, "Source: Market Intelligence Report 2024")


    # ─────────────────────────────────────────────────────────────────────────
    # Slide 7: 2×2 Priority Matrix
    # ─────────────────────────────────────────────────────────────────────────
    def slide7():
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        slide = blank_slide()
        add_title(slide, "Initiative Priority Matrix", "Implementation Ease vs. Business Impact")

        # Chart canvas
        CX = Inches(0.6)
        CY = Inches(1.4)
        CW = Inches(11.8)
        CH = Inches(5.3)

        # Quadrant backgrounds
        quad_colors = [
            ("#EAF4FB", CX,        CY,        CW/2, CH/2),   # top-left  = Hard/High
            ("#EBF7E6", CX+CW/2,   CY,        CW/2, CH/2),   # top-right = Easy/High  *
            ("#F5F5F5", CX,        CY+CH/2,   CW/2, CH/2),   # bot-left  = Hard/Low
            ("#FEF9E7", CX+CW/2,   CY+CH/2,   CW/2, CH/2),   # bot-right = Easy/Low
        ]
        quad_labels = [
            ("Strategic\nBets",   CX+CW/4,     CY+CH/4),
            ("Quick Wins *",      CX+3*CW/4,   CY+CH/4),
            ("Avoid / Defer",     CX+CW/4,     CY+3*CH/4),
            ("Fill-ins",          CX+3*CW/4,   CY+3*CH/4),
        ]
        for color, qx, qy, qw, qh in quad_colors:
            rect = slide.shapes.add_shape(1, int(qx), int(qy), int(qw), int(qh))
            rect.fill.solid(); rect.fill.fore_color.rgb = _hex(color)
            rect.line.color.rgb = _hex(LIGHT_GRAY); rect.line.width = Pt(0.5)

        for qlbl, qx, qy in quad_labels:
            txb = slide.shapes.add_textbox(int(qx-Inches(0.8)), int(qy-Inches(0.2)), int(Inches(1.6)), int(Inches(0.4)))
            p = txb.text_frame.paragraphs[0]
            p.text = qlbl; p.font.size = Pt(9); p.font.name = "Arial"
            p.font.color.rgb = _hex(MID_GRAY); p.font.italic = True
            p.alignment = PP_ALIGN.CENTER

        # Midpoint axis lines
        for conn_coords in [
            (int(CX), int(CY+CH/2), int(CX+CW), int(CY+CH/2)),
            (int(CX+CW/2), int(CY), int(CX+CW/2), int(CY+CH)),
        ]:
            ln = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, *conn_coords)
            ln.line.color.rgb = _hex(MID_GRAY); ln.line.width = Pt(0.75)

        # Plot dots + labels  (ease 0-10 -> x, impact 0-10 -> y)
        def px(ease):   return int(CX + CW * ease / 10)
        def py(impact): return int(CY + CH * (1 - impact / 10))

        colors_by_quadrant = {
            (True, True):   DARK_BLUE,
            (True, False):  MID_BLUE,
            (False, True):  MID_BLUE,
            (False, False): MID_GRAY,
        }
        for label_txt, ease, impact in MATRIX_ITEMS:
            easy  = ease   > 5
            high  = impact > 5
            color = colors_by_quadrant[(easy, high)]
            cx_, cy_ = px(ease), py(impact)
            dot = slide.shapes.add_shape(9, cx_-int(Inches(0.15)), cy_-int(Inches(0.15)),
                                          int(Inches(0.3)), int(Inches(0.3)))
            dot.fill.solid(); dot.fill.fore_color.rgb = _hex(color)
            dot.line.color.rgb = _hex(WHITE); dot.line.width = Pt(0.5)

            # Label offset: prefer right, flip if near edge
            lx = cx_ + int(Inches(0.18))
            if ease > 8.5: lx = cx_ - int(Inches(1.4))
            ly = cy_ - int(Inches(0.2))
            txb = slide.shapes.add_textbox(lx, ly, int(Inches(1.2)), int(Inches(0.4)))
            p = txb.text_frame.paragraphs[0]
            p.text = label_txt.replace("\n", " ")
            p.font.size = Pt(8); p.font.name = "Arial"
            p.font.color.rgb = _hex(color); p.font.bold = True

        # Axis labels
        for lbl_txt, x, y in [
            ("<- Harder  |  Implementation Ease  |  Easier ->", int(CX+CW/2), int(CY+CH+Inches(0.15))),
            ("Low Impact", int(CX-Inches(0.5)), int(CY+3*CH/4)),
            ("High Impact", int(CX-Inches(0.55)), int(CY+CH/4)),
        ]:
            txb = slide.shapes.add_textbox(int(x-Inches(2)), int(y-Inches(0.2)), int(Inches(4)), int(Inches(0.35)))
            p = txb.text_frame.paragraphs[0]
            p.text = lbl_txt; p.font.size = Pt(9); p.font.name = "Arial"
            p.font.color.rgb = _hex(DARK_GRAY); p.alignment = PP_ALIGN.CENTER

        add_source(slide, "Source: Strategic Planning Workshop 2024")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 8: EBITDA Bridge (horizontal waterfall)
    # ─────────────────────────────────────────────────────────────────────────
    def slide8():
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        slide = blank_slide()
        add_title(slide, "EBITDA Bridge", "$M")

        # Canvas (horizontal - values run left to right)
        CX   = Inches(1.5)
        CY   = Inches(1.5)
        CW   = Inches(10.8)
        CH   = Inches(4.8)
        XMAX = 550   # x-axis max (slightly above 500)
        BAR_H_EACH = int(CH / (len(EBITDA_ITEMS) + 0.5))
        BH = int(BAR_H_EACH * 0.65)
        GAP = int(BAR_H_EACH * 0.175)

        def xc(v): return int(CX + CW * v / XMAX)

        COLOR_START = DARK_BLUE
        COLOR_UP    = "#70AD47"
        COLOR_DOWN  = "#C00000"
        COLOR_TOTAL = DARK_BLUE

        running = 0
        for i, (lbl_txt, val, kind) in enumerate(EBITDA_ITEMS):
            by = int(CY + GAP + i * BAR_H_EACH)
            if kind == "start":
                bx = xc(0); bw = xc(val) - xc(0)
                color = COLOR_START; running = val
            elif kind == "total":
                bx = xc(0); bw = xc(running) - xc(0)
                color = COLOR_TOTAL
            elif kind == "up":
                bx = xc(running); bw = xc(running + val) - xc(running)
                color = COLOR_UP; running += val
            else:  # down
                new_run = running + val  # val is negative
                bx = xc(new_run); bw = xc(running) - xc(new_run)
                color = COLOR_DOWN; running = new_run

            rect = slide.shapes.add_shape(1, bx, by, bw, BH)
            rect.fill.solid(); rect.fill.fore_color.rgb = _hex(color)
            rect.line.fill.background()

            # Value label inside or to right
            lbl_x = bx + bw + int(Inches(0.08))
            lbl_y = by
            txb = slide.shapes.add_textbox(lbl_x, lbl_y, int(Inches(0.7)), BH)
            p = txb.text_frame.paragraphs[0]
            p.text = f"${abs(val)}M"
            p.font.size = Pt(9); p.font.bold = True; p.font.name = "Arial"
            p.font.color.rgb = _hex(color if kind != "total" else DARK_BLUE)

            # Row label on left
            ltxb = slide.shapes.add_textbox(int(CX - Inches(1.4)), by, int(Inches(1.35)), BH)
            ltxb.text_frame.word_wrap = True
            lp = ltxb.text_frame.paragraphs[0]
            lp.text = lbl_txt; lp.font.size = Pt(9); lp.font.name = "Arial"
            lp.font.color.rgb = _hex(DARK_GRAY); lp.alignment = PP_ALIGN.RIGHT

        # X-axis line
        ax_y = CY + CH
        ax = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, int(CX), int(ax_y), int(CX+CW), int(ax_y))
        ax.line.color.rgb = _hex(MID_GRAY); ax.line.width = Pt(0.75)

        # X-axis gridlines every 100
        for gv in range(0, 600, 100):
            gx = xc(gv)
            gl = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, gx, int(CY), gx, int(ax_y))
            gl.line.color.rgb = _hex(LIGHT_GRAY); gl.line.width = Pt(0.5)
            txb = slide.shapes.add_textbox(gx - int(Inches(0.3)), int(ax_y + Inches(0.05)),
                                            int(Inches(0.6)), int(Inches(0.25)))
            p = txb.text_frame.paragraphs[0]
            p.text = f"${gv}M"; p.font.size = Pt(8); p.font.name = "Arial"
            p.font.color.rgb = _hex(MID_GRAY); p.alignment = PP_ALIGN.CENTER

        add_source(slide, "Source: Management Accounts FY2024")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 9: Gantt / Transformation Roadmap
    # ─────────────────────────────────────────────────────────────────────────
    def slide9():
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        slide = blank_slide()
        add_title(slide, "Transformation Roadmap", "Weeks")

        CX = Inches(2.8)
        CY = Inches(1.4)
        CW = Inches(9.8)
        CH = Inches(5.2)
        N  = len(GANTT_DATA)
        ROW_H = int(CH / (N + 0.5))
        BH    = int(ROW_H * 0.62)
        GAP   = int(ROW_H * 0.19)

        def gx(week): return int(CX + CW * week / GANTT_TOTAL_WEEKS)

        # Phase dividers / header
        for q in range(0, GANTT_TOTAL_WEEKS + 1, 4):
            lx = gx(q)
            gl = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, lx, int(CY), lx, int(CY+CH))
            gl.line.color.rgb = _hex(LIGHT_GRAY); gl.line.width = Pt(0.5)
            if q % 8 == 0:
                txb = slide.shapes.add_textbox(lx - int(Inches(0.3)), int(CY - Inches(0.35)),
                                                int(Inches(0.7)), int(Inches(0.3)))
                p = txb.text_frame.paragraphs[0]
                p.text = f"W{q}"; p.font.size = Pt(8); p.font.name = "Arial"
                p.font.color.rgb = _hex(MID_GRAY); p.alignment = PP_ALIGN.CENTER

        # Draw task bars
        prev_ws = None
        for i, (ws_name, task, start, dur, color) in enumerate(GANTT_DATA):
            by = int(CY + GAP + i * ROW_H)
            bx = gx(start)
            bw = gx(start + dur) - gx(start)
            rect = slide.shapes.add_shape(1, bx, by, bw, BH)
            rect.fill.solid(); rect.fill.fore_color.rgb = _hex(color)
            rect.line.fill.background()
            # Task label inside bar if wide enough
            if bw > int(Inches(0.8)):
                txb = slide.shapes.add_textbox(bx + int(Inches(0.07)), by, bw - int(Inches(0.1)), BH)
                txb.text_frame.word_wrap = False
                p = txb.text_frame.paragraphs[0]
                p.text = task; p.font.size = Pt(8); p.font.name = "Arial"
                p.font.color.rgb = _hex(WHITE if color != "#D9D9D9" else DARK_GRAY)
            # Workstream label on left  (only first row of each ws)
            if ws_name != prev_ws:
                ltxb = slide.shapes.add_textbox(int(CX - Inches(2.7)), by, int(Inches(2.6)), BH)
                lp = ltxb.text_frame.paragraphs[0]
                lp.text = ws_name; lp.font.size = Pt(9.5); lp.font.bold = True
                lp.font.name = "Arial"; lp.font.color.rgb = _hex(DARK_BLUE)
                lp.alignment = PP_ALIGN.RIGHT
            prev_ws = ws_name

        # X baseline
        ax = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, int(CX), int(CY+CH), int(CX+CW), int(CY+CH))
        ax.line.color.rgb = _hex(MID_GRAY); ax.line.width = Pt(0.75)

        add_source(slide, "Source: Programme Management Office")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 10: Harvey Balls - Capability Benchmarking
    # ─────────────────────────────────────────────────────────────────────────
    def slide10():
        import math
        slide = blank_slide()
        add_title(slide, "Capability Benchmarking", "Harvey Ball Assessment")

        CELL_W = int(Inches(1.9))
        CELL_H = int(Inches(0.62))
        COL0_W = int(Inches(2.8))
        START_X = int(Inches(0.5))
        START_Y = int(Inches(1.45))
        BALL_R  = int(Inches(0.22))

        # Header row
        hdr_bg = slide.shapes.add_shape(1, START_X, START_Y, COL0_W + CELL_W * len(HARVEY_COMPANIES), CELL_H)
        hdr_bg.fill.solid(); hdr_bg.fill.fore_color.rgb = _hex(DARK_BLUE)
        hdr_bg.line.fill.background()
        for c, co in enumerate(HARVEY_COMPANIES):
            cx = START_X + COL0_W + c * CELL_W
            txb = slide.shapes.add_textbox(cx, START_Y, CELL_W, CELL_H)
            p = txb.text_frame.paragraphs[0]
            p.text = co; p.font.size = Pt(10); p.font.bold = True
            p.font.name = "Arial"; p.font.color.rgb = _hex(WHITE)
            p.alignment = PP_ALIGN.CENTER

        def draw_harvey(cx_center, cy_center, score, radius):
            """Draw a Harvey Ball: score 0-4."""
            # Background circle (white fill, dark border)
            r2 = radius * 2
            bg = slide.shapes.add_shape(9,  # oval
                cx_center - radius, cy_center - radius, r2, r2)
            bg.fill.solid(); bg.fill.fore_color.rgb = _hex(WHITE)
            bg.line.color.rgb = _hex(DARK_BLUE); bg.line.width = Pt(1.5)
            if score == 0:
                return
            if score == 4:
                # Fully filled
                dot = slide.shapes.add_shape(9,
                    cx_center - radius, cy_center - radius, r2, r2)
                dot.fill.solid(); dot.fill.fore_color.rgb = _hex(DARK_BLUE)
                dot.line.fill.background()
                return
            # Partial fill: use a freeform pie-slice as polygon
            # build_freeform uses EMU coordinates directly (scale=1.0)
            frac = score / 4.0
            n_pts = max(16, int(frac * 48))
            angles = [(-math.pi/2) + frac * 2 * math.pi * k / n_pts for k in range(n_pts + 1)]
            arc_pts = [(cx_center + int(radius * math.cos(a)),
                        cy_center + int(radius * math.sin(a))) for a in angles]
            # Polygon: center -> arc points -> back to center (closed)
            all_pts = [(cx_center, cy_center)] + arc_pts + [(cx_center, cy_center)]
            builder = slide.shapes.build_freeform(all_pts[0][0], all_pts[0][1])
            builder.add_line_segments(all_pts[1:], close=True)
            shape = builder.convert_to_shape()
            shape.fill.solid(); shape.fill.fore_color.rgb = _hex(DARK_BLUE)
            shape.line.fill.background()

        for r, (crit, row_scores) in enumerate(zip(HARVEY_CRITERIA, HARVEY_SCORES)):
            ry = START_Y + CELL_H + r * CELL_H
            # Row bg
            bg_color = WHITE if r % 2 == 0 else "#F5F8FF"
            row_bg = slide.shapes.add_shape(1, START_X, ry, COL0_W + CELL_W * len(HARVEY_COMPANIES), CELL_H)
            row_bg.fill.solid(); row_bg.fill.fore_color.rgb = _hex(bg_color)
            row_bg.line.color.rgb = _hex(LIGHT_GRAY); row_bg.line.width = Pt(0.25)
            # Criterion label
            txb = slide.shapes.add_textbox(START_X + int(Inches(0.1)), ry, COL0_W - int(Inches(0.15)), CELL_H)
            p = txb.text_frame.paragraphs[0]
            p.text = crit; p.font.size = Pt(9.5); p.font.name = "Arial"
            p.font.color.rgb = _hex(DARK_GRAY)
            # Harvey balls
            for c, score in enumerate(row_scores):
                cx_center = START_X + COL0_W + c * CELL_W + CELL_W // 2
                cy_center = ry + CELL_H // 2
                draw_harvey(cx_center, cy_center, score, BALL_R)

        # Legend
        ly = START_Y + CELL_H + len(HARVEY_CRITERIA) * CELL_H + int(Inches(0.2))
        for i, desc in enumerate(["No capability", "Basic", "Moderate", "Strong", "Leading"]):
            lx = START_X + i * int(Inches(2.4))
            draw_harvey(lx + int(Inches(0.25)), ly + int(Inches(0.22)), i, int(Inches(0.17)))
            txb = slide.shapes.add_textbox(lx + int(Inches(0.52)), ly, int(Inches(1.8)), int(Inches(0.44)))
            p = txb.text_frame.paragraphs[0]
            p.text = desc; p.font.size = Pt(8.5); p.font.name = "Arial"
            p.font.color.rgb = _hex(DARK_GRAY)

        add_source(slide, "Source: Internal capability assessment")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 11: Tornado / Sensitivity Analysis
    # ─────────────────────────────────────────────────────────────────────────
    def slide11():
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        slide = blank_slide()
        add_title(slide, "NPV Sensitivity Analysis", "Impact on base-case NPV of $217M  ($M)")

        sorted_drivers = sorted(TORNADO_DRIVERS, key=lambda x: abs(x[1])+abs(x[2]), reverse=True)
        N = len(sorted_drivers)

        CX   = Inches(3.2)
        CY   = Inches(1.45)
        CW   = Inches(9.3)
        CH   = Inches(5.1)
        XMIN = -60
        XMAX =  60
        XRNG = XMAX - XMIN

        ROW_H = int(CH / (N + 0.5))
        BH    = int(ROW_H * 0.62)
        GAP   = int(ROW_H * 0.19)

        def xc(v): return int(CX + CW * (v - XMIN) / XRNG)

        # Zero line
        x0 = xc(0)
        zl = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x0, int(CY), x0, int(CY+CH))
        zl.line.color.rgb = _hex(DARK_GRAY); zl.line.width = Pt(1)

        # Gridlines
        for gv in range(-60, 61, 20):
            gx = xc(gv)
            gl = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, gx, int(CY), gx, int(CY+CH))
            gl.line.color.rgb = _hex(LIGHT_GRAY); gl.line.width = Pt(0.5)
            txb = slide.shapes.add_textbox(gx - int(Inches(0.3)), int(CY+CH+Inches(0.05)),
                                            int(Inches(0.6)), int(Inches(0.25)))
            p = txb.text_frame.paragraphs[0]
            p.text = f"${gv}M" if gv != 0 else "Base"
            p.font.size = Pt(8); p.font.name = "Arial"
            p.font.color.rgb = _hex(MID_GRAY); p.alignment = PP_ALIGN.CENTER

        for i, (drv, low, high) in enumerate(sorted_drivers):
            by = int(CY + GAP + i * ROW_H)

            # Downside bar (negative)
            bx_low = xc(low); bx_0 = xc(0)
            bw_low = bx_0 - bx_low
            r_down = slide.shapes.add_shape(1, bx_low, by, bw_low, BH)
            r_down.fill.solid(); r_down.fill.fore_color.rgb = _hex(RED)
            r_down.line.fill.background()

            # Upside bar (positive)
            bx_hi = xc(0); bw_hi = xc(high) - bx_hi
            r_up = slide.shapes.add_shape(1, bx_hi, by, bw_hi, BH)
            r_up.fill.solid(); r_up.fill.fore_color.rgb = _hex(GREEN)
            r_up.line.fill.background()

            # Driver label on left
            ltxb = slide.shapes.add_textbox(int(CX - Inches(3.0)), by, int(Inches(2.85)), BH)
            lp = ltxb.text_frame.paragraphs[0]
            lp.text = drv; lp.font.size = Pt(9); lp.font.name = "Arial"
            lp.font.color.rgb = _hex(DARK_GRAY); lp.alignment = PP_ALIGN.RIGHT

            # Value labels
            for bx_lbl, bw_lbl, val, align_right in [
                (bx_low - int(Inches(0.55)), int(Inches(0.52)), low, True),
                (xc(high) + int(Inches(0.05)), int(Inches(0.52)), high, False),
            ]:
                txb = slide.shapes.add_textbox(bx_lbl, by, int(Inches(0.52)), BH)
                p = txb.text_frame.paragraphs[0]
                p.text = f"${val:+d}M"; p.font.size = Pt(8); p.font.name = "Arial"
                p.font.color.rgb = _hex(RED if val < 0 else GREEN); p.font.bold = True
                p.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT

        add_source(slide, "Source: Financial model - sensitivity run")

    # ─────────────────────────────────────────────────────────────────────────
    # Slide 12: Sankey - Income Statement
    # ─────────────────────────────────────────────────────────────────────────
    def slide12():
        slide = blank_slide()
        add_title(slide, "Income Statement Sankey", "FY 2024  ($B)")
        png = _make_sankey_png()
        if png:
            slide.shapes.add_picture(png, Inches(0.25), Inches(1.1), Inches(12.8), Inches(6.2))
        else:
            txb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
            p = txb.text_frame.paragraphs[0]
            p.text = "Install matplotlib to render Sankey:  pip install matplotlib"
            p.font.size = Pt(14); p.font.name = "Arial"; p.font.color.rgb = _hex(MID_GRAY)
        add_source(slide, "Source: Company financials FY2024")

    slide1(); slide2(); slide3(); slide4(); slide5(); slide6(); slide7(); slide8(); slide9(); slide10(); slide11(); slide12()
    prs.save("McKinsey_Slides.pptx")
    print("  ✓  McKinsey_Slides.pptx   created")


# ══════════════════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("\nMcKinsey / BCG Consulting Chart Generator")
    print("=" * 42)
    build_excel()
    build_pptx()
    print("\nDone! Open the files in Excel / PowerPoint and replace the sample data.")
    print("Each chart updates automatically when you change the values in the data table.")
